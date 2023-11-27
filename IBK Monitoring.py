import json
import requests
import matplotlib.pyplot as plt
from io import BytesIO
from datetime import datetime, timedelta
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os
from tqdm import tqdm

def get_auth_token():
    script_path = os.path.dirname(os.path.realpath(__file__))

    config_path = os.path.join(script_path, 'config.json')

    try:
        with open(config_path, 'r') as config_file:
            config = json.load(config_file)

        url = f"https://iam.{config['endpoint']}.myhuaweicloud.com/v3/auth/tokens"

        payload = {
            "auth": {
                "identity": {
                    "methods": ["password"],
                    "password": {
                        "user": {
                            "name": config["username"],
                            "password": config["password"],
                            "domain": {"name": config["tenant_name"]},
                        }
                    }
                },
                "scope": {"project": {"name": config["endpoint"]}},
            }
        }

        headers = {"Content-Type": "application/json"}

        response = requests.post(url, headers=headers, json=payload)

        if response.ok:
            subject_token = response.headers.get('X-Subject-Token')

            config["auth_token"] = subject_token

            with open(config_path, 'w') as config_file:
                json.dump(config, config_file, indent=2)

            print("Token actualizado en el archivo config.json.")
        else:
            print(f"Error: {response.status_code}, {response.text}")

    except Exception as e:
        print(f"Error: {e}")

def get_ecs_list():
    # Obtener la ruta del script
    script_path = os.path.dirname(os.path.realpath(__file__))

    # Construir la ruta completa del archivo de configuración
    config_path = os.path.join(script_path, 'config.json')

    try:
        # Leer la configuración desde el archivo
        with open(config_path, 'r') as config_file:
            config = json.load(config_file)

        # Construir la URL de la solicitud
        url = f"https://ecs.{config['endpoint']}.myhuaweicloud.com/v1/{config['project_id']}/cloudservers/detail?offset=1&status=ACTIVE"

        # Encabezados de la solicitud
        headers = {
            'X-Auth-Token': config['auth_token'],
            'Cookie': 'HWWAFSESID=00ffe73528d310d20d; HWWAFSESTIME=1700877930137'
        }

        # Realizar la solicitud GET
        response = requests.get(url, headers=headers)

        # Verificar si la solicitud fue exitosa (código de estado 2xx)
        if response.ok:
            # Convertir la respuesta JSON a un diccionario
            response_json = response.json()

            # Verificar si la clave "servers" está presente en la respuesta
            if 'servers' in response_json:
                # Crear una lista con los valores de "id" y "name" para cada objeto en "servers"
                server_list = [{'id': server['id'], 'name': server['name']} for server in response_json['servers']]

                # Retornar la nueva lista
                return server_list
            else:
                print("La clave 'servers' no está presente en la respuesta.")

        else:
            # Imprimir el código de estado y el mensaje de error
            print(f"Error: {response.status_code}, {response.text}")

    except Exception as e:
        print(f"Error: {e}")

    # En caso de error, retornar una lista vacía
    return []

def hex_to_rgb(hex_color):
    return RGBColor(int(hex_color[1:3], 16), int(hex_color[3:5], 16), int(hex_color[5:7], 16))

def move_slides_to_end(prs, slide_numbers):
    sld_id_lst = prs.slides._sldIdLst
    slides_to_move = [sld_id_lst[s - 1] for s in slide_numbers]

    for slide_id in slides_to_move:
        sld_id_lst.remove(slide_id)
        sld_id_lst.insert(len(sld_id_lst), slide_id)

def add_slide_with_text_and_images(presentation, title, font_color, font_name, font_size, width, height, left, top, subtitle_text=None, subtitle_color=None, subtitle_size=None, img_paths=None):
    slide_layout = presentation.slide_layouts[6]
    slide = presentation.slides.add_slide(slide_layout)

    text_box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    text_frame = text_box.text_frame

    p = text_frame.add_paragraph()
    p.text = title
    p.font.name = font_name
    p.font.size = Pt(int(font_size))
    p.font.color.rgb = hex_to_rgb(font_color)
    p.font.bold = True
    p.font.character_spacing = Pt(2.0)
    p.alignment = PP_ALIGN.CENTER
    p.space_before = Pt(0)
    p.space_after = Pt(0)

    text_frame.word_wrap = True
    text_frame.auto_size = True
    text_frame.margin_top = Pt(0)
    text_frame.margin_bottom = Pt(0)

    text_box.left = int((presentation.slide_width - Inches(width)) / 2)

    if subtitle_text:
        subtitle_box = slide.shapes.add_textbox(Inches(left), Inches(top + height), Inches(width), Inches(0.5))
        subtitle_frame = subtitle_box.text_frame

        subtitle_paragraph = subtitle_frame.add_paragraph()
        subtitle_paragraph.text = subtitle_text
        subtitle_paragraph.font.name = font_name
        subtitle_paragraph.font.size = Pt(12)
        subtitle_paragraph.font.color.rgb = hex_to_rgb(subtitle_color) if subtitle_color else hex_to_rgb("#44546A")
        subtitle_paragraph.font.bold = True
        subtitle_paragraph.alignment = PP_ALIGN.CENTER
        subtitle_paragraph.space_before = Pt(0)
        subtitle_paragraph.space_after = Pt(0)

        subtitle_frame.word_wrap = True
        subtitle_frame.auto_size = True

    if img_paths:
        # Cambiar las ubicaciones de las imágenes según el nuevo diseño
        img_locations = [(0.67, 1.74), (8.67, 1.74), (4.66, 4.25)]  # CPU, RAM, Disk
        for img_path, img_location in zip(img_paths, img_locations):
            left_offset, top_offset = img_location
            img_left = Inches(left_offset)
            img_top = Inches(top_offset)
            img_width = Inches(4)
            img_height = Inches(3)

            if os.path.exists(img_path):
                slide.shapes.add_picture(img_path, img_left, img_top, width=img_width, height=img_height)

def create_folder_if_not_exists(folder_path):
    if not os.path.exists(folder_path):
        os.makedirs(folder_path)

def get_metric_data(instance_id, metric_name):
    # Obtener la ruta del script actual
    script_path = os.path.dirname(os.path.realpath(__file__))
    
    # Combinar la ruta del script con el nombre del archivo de configuración
    config_path = os.path.join(script_path, 'config.json')

    # Cargar valores desde el archivo de configuración
    with open(config_path, 'r') as config_file:
        config = json.load(config_file)

    to_timestamp = int(datetime.utcnow().timestamp() * 1000)
    from_timestamp = int((datetime.utcnow() - timedelta(days=7)).timestamp() * 1000)
    
    url = f"https://ces.{config['endpoint']}.myhuaweicloud.com/V1.0/{config['project_id']}/metric-data?namespace=SYS.ECS&metric_name={metric_name}&dim.0=instance_id,{instance_id}&from={from_timestamp}&to={to_timestamp}&period=3600&filter=max"

    headers = {
        'X-Auth-Token': config['auth_token']
    }

    response = requests.get(url, headers=headers)
    data = response.json()

    return data

def convert_timestamp_to_date(timestamp):
    timestamp_in_seconds = timestamp / 1000
    return datetime.utcfromtimestamp(timestamp_in_seconds).strftime('%d %b - %H:%M')

def plot_and_save_graph(instance_id, instance_name, metric_name, folder_path, title):
    try:
        data = get_metric_data(instance_id, metric_name)
        datapoints = data['datapoints']

        values = [entry['max'] for entry in datapoints]
        dates = [convert_timestamp_to_date(entry['timestamp']) for entry in datapoints]

        max_value_index = values.index(max(values))
        max_value_date = convert_timestamp_to_date(datapoints[max_value_index]['timestamp'])

        plt.plot(dates, values, marker='o', linestyle='-', label=f'{metric_name.capitalize()} Utilization')
        plt.scatter(max_value_date, max(values), color='red', zorder=5)
        plt.annotate(
            f'{convert_timestamp_to_date(datapoints[max_value_index]["timestamp"])}\nMax: {max(values)}%',
            (max_value_date, max(values)), textcoords="offset points", xytext=(0, -40), ha='center',
            bbox=dict(boxstyle='round,pad=0.3', facecolor='white', edgecolor='black')
        )

        plt.xlabel('Línea de tiempo')
        plt.ylabel(f'{title} (%)')
        plt.title(f'{title}' + " de " + f'{instance_name}')
        plt.xticks([])

        plt.grid(True)
        plt.tight_layout()

        img_path = os.path.join(folder_path, f'{metric_name}_graph.png')
        plt.savefig(img_path)
        plt.close()

    except Exception as e:
        print(f"Error al obtener datos para {metric_name} de la instancia {instance_name}: {e}")

def make_ppt(ecs_list):
    img_progress_bar = tqdm(total=len(ecs_list), desc="Generando PPT")

    prs = Presentation("template.pptx")

    for ecs_instance in ecs_list:
        folder_path = os.path.join("Imágenes", ecs_instance["name"])
        create_folder_if_not_exists(folder_path)

        plot_and_save_graph(ecs_instance["id"], ecs_instance["name"], "cpu_util", folder_path, "Uso de CPU")
        plot_and_save_graph(ecs_instance["id"], ecs_instance["name"], "mem_util", folder_path, "Uso de RAM")
        plot_and_save_graph(ecs_instance["id"], ecs_instance["name"], "disk_util_inband", folder_path, "Uso de disco")

        img_progress_bar.update(1)

        add_slide_with_text_and_images(
            prs,
            'Reportes de operación',
            "#263E83",
            "Montserrat",
            40,
            7.47,
            0.72,
            2.83,
            0.37,
            subtitle_text=f"Terminal: {ecs_instance['name']}",
            subtitle_color="#44546A",
            subtitle_size=12,
            img_paths=[
                os.path.join("Imágenes", ecs_instance["name"], 'cpu_util_graph.png'),
                os.path.join("Imágenes", ecs_instance["name"], 'mem_util_graph.png'),
                os.path.join("Imágenes", ecs_instance["name"], 'disk_util_inband_graph.png'),
            ]
        )

    img_progress_bar.close()

    # Especificar los números de las diapositivas que deseas mover al final
    slides_to_move_numbers = [5, 6, 7]

    # Mover las diapositivas al final
    move_slides_to_end(prs, slides_to_move_numbers)

    prs.save('Monitoreo.pptx')

def main():
    get_auth_token()
    ecs_list = get_ecs_list()
    print(len(ecs_list))
    make_ppt(ecs_list)



if __name__ == "__main__":
    main()
